from pptx import Presentation
from pptx.util import Inches

# Создание презентации
prs = Presentation()

# Функция для добавления слайда с заголовком и содержанием
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Слайд с заголовком и содержанием
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

# Слайды презентации
slides = [
    ("Введение в API и Postman", "Урок 1\nАвтор: [Ваше Имя]"),
    ("Приветствие и представление", "Приветствие зрителей\nКраткое представление\nЦель курса\nОписание структуры урока"),
    ("Что такое API?", "Определение API\nПростое объяснение API как посредника\nПримеры: погода, карты, соцсети"),
    ("Зачем нужны API?", "Важность API для интеграции систем\nПримеры использования API\nРоль API в автоматизации тестирования и разработке"),
    ("Популярные виды API", "REST API:\nПринципы REST\nПримеры использования\n\nSOAP API:\nПротокол обмена сообщениями\nПримеры использования\n\nGraphQL API:\nЗапросы с одной конечной точки\nПримеры использования"),
    ("REST API", "Representational State Transfer\nИспользование HTTP-запросов\nПримеры: GET, POST, PUT, DELETE"),
    ("SOAP API", "Simple Object Access Protocol\nИспользование XML\nНадежность и безопасность"),
    ("GraphQL API", "Запросы с одной конечной точки\nЗапрос только необходимых данных\nОптимизация работы с данными"),
    ("Поддержка API в Postman", "Поддержка различных типов API\nСоздание и отправка запросов\nАнализ ответов\nАвтоматизация тестирования"),
    ("Заключение", "Краткое резюме\nАнонс следующего урока (обзор интерфейса Postman)\nПризыв к подписке, лайкам и комментариям")
]

# Добавление слайдов в презентацию
for title, content in slides:
    add_slide(title, content)

# Сохранение презентации
file_path = "./Введение в API и Postman.pptx"
prs.save(file_path)

file_path
